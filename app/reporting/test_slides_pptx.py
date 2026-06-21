# -*- coding: utf-8 -*-
"""PPTX 슬라이드 보고서 빌더 테스트."""
import io
import unittest

from app.reporting.slides_pptx import build_slides_pptx

_META = {"name": "지니스안경", "campaign_name": "4월 누진렌즈", "region": "공주", "period": "26.4.1~30"}
_KPI = {
    "total_cost": 328095, "total_impressions": 103010, "total_clicks": 1373,
    "ctr": 1.33, "cpc": 239, "cpm": 3185, "total_inquiries": 9, "cpa": 36455,
    "total_regulars": 96, "cpr": 3418, "total_coupons": 112, "cp_coupon": 2929,
}
_INSIGHTS = {
    "conclusion": "노출·클릭 정상, 클릭→문의 병목.",
    "good": "노출·클릭 단계 정상", "blocked": "클릭→문의 전환 낮음",
    "next_actions": ["소식글 CTA를 쿠폰 먼저로 재정비", "30~40대 여성 타겟 축소"],
    "experiments": [
        {"priority": "1", "change": "CTA 쿠폰 먼저", "success_criteria": "쿠폰 전환↑", "owner": "사장님", "schedule": "이번 주"},
    ],
    "judgment": {"expand": "CPA 기준 이하면 증액", "review": "기준 근처면 소재 재정의", "stop": "정체 시 중단"},
}
_FUNNEL = [
    {"label": "노출", "count": 103010, "cost_per": 3185, "cost_label": "1천회당"},
    {"label": "클릭", "count": 1373, "rate_label": "노출 대비", "rate": 1.33, "cost_per": 239, "cost_label": "1회당"},
]


class TestBuildSlidesPptx(unittest.TestCase):
    def test_returns_pptx_bytes(self):
        data = build_slides_pptx(_META, _KPI, _INSIGHTS, funnel_stages=_FUNNEL)
        self.assertIsInstance(data, bytes)
        self.assertGreater(len(data), 1000)
        # PPTX = zip(PK 시그니처)
        self.assertEqual(data[:2], b"PK")

    def test_reopens_with_expected_slides(self):
        from pptx import Presentation
        seg = [{"label": "30-40 여성", "cost": 120000, "cpa": 9000, "verdict": "증액"}]
        data = build_slides_pptx(_META, _KPI, _INSIGHTS, funnel_stages=_FUNNEL, segment_rows=seg, generated_at="2026-06-20")
        prs = Presentation(io.BytesIO(data))
        # 표지·지표·퍼널·진단·세그먼트·액션·판단 = 7
        self.assertEqual(len(prs.slides), 7)

    def test_minimal_inputs_no_crash(self):
        data = build_slides_pptx(_META, _KPI, {})
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        # 표지 + 지표만 (인사이트/퍼널 없음)
        self.assertEqual(len(prs.slides), 2)


if __name__ == "__main__":
    unittest.main()
