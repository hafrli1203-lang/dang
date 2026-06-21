# -*- coding: utf-8 -*-
"""성과+세그먼트 통합 보고서를 편집 가능한 PowerPoint(.pptx)로.

HTML 슬라이드(slides_html.py)와 동일 구성을 파워포인트로 — 광고주에게 보내고
브랜딩/문구를 직접 고칠 수 있게. build_slides_pptx()는 bytes를 반환(순수, 파일 I/O 없음).
python-pptx 필요(requirements.txt). 폰트는 Paperlogy 우선, 없으면 맑은 고딕.
"""
from __future__ import annotations

import io
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_ORANGE = RGBColor(0xFF, 0x6F, 0x0F)
_INK = RGBColor(0x1A, 0x1A, 0x2E)
_MUTED = RGBColor(0x6B, 0x72, 0x80)
_LINE = RGBColor(0xE7, 0xE2, 0xDA)
_BG = RGBColor(0xF8, 0xF9, 0xFC)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_FONT = "Paperlogy"
_FONT_FALLBACK = "맑은 고딕"

_W = Inches(13.333)
_H = Inches(7.5)


def _won(v) -> str:
    try:
        return f"{int(round(float(v))):,}원"
    except (TypeError, ValueError):
        return "-"


def _num(v) -> str:
    try:
        return f"{int(round(float(v))):,}"
    except (TypeError, ValueError):
        return "-"


def _pct(v) -> str:
    try:
        return f"{float(v):.1f}%"
    except (TypeError, ValueError):
        return "-"


def _set_font(run, size, bold=False, color=_INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = _FONT


def _textbox(slide, x, y, w, h, lines, *, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, bold, color, align)] — 여러 단락."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_font(run, size, bold, color)
    return tb


def _card(slide, x, y, w, h, fill=_BG, line=_LINE, top_accent=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if top_accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _ORANGE
        bar.line.fill.background()
        bar.shadow.inherit = False
    return shp


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text):
    _textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9),
             [(text, 30, True, _INK, PP_ALIGN.LEFT)])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Pt(6), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ORANGE
    bar.line.fill.background()
    bar.shadow.inherit = False


def _table(slide, x, y, w, headers, rows, col_widths=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    h = Inches(0.45 * nrows)
    gtbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            gtbl.columns[ci].width = cw
    for ci, head in enumerate(headers):
        cell = gtbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _BG
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(head)
        _set_font(run, 13, True, _INK)
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _WHITE
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            _set_font(run, 12, False, _INK)
    return gtbl


def build_slides_pptx(
    meta: dict,
    kpi: dict,
    insights: dict,
    *,
    funnel_stages: Optional[list] = None,
    segment_rows: Optional[list] = None,
    generated_at: str = "",
) -> bytes:
    """통합 성과 보고서 PPTX(bytes). 인자는 slides_html.build_slides_html과 동일."""
    meta = meta or {}
    kpi = kpi or {}
    insights = insights or {}
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    # 1) 표지
    s = _blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _W, _H)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xEC)
    bg.line.fill.background(); bg.shadow.inherit = False
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.6), Inches(2.6), Inches(0.5))
    pill.fill.solid(); pill.fill.fore_color.rgb = _ORANGE
    pill.line.fill.background(); pill.shadow.inherit = False
    ptf = pill.text_frame; ptf.word_wrap = True
    pr = ptf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    prr = pr.add_run(); prr.text = "당근 광고 성과 보고서"; _set_font(prr, 13, True, _WHITE)
    sub = " · ".join(b for b in [
        str(meta.get("campaign_name") or ""), str(meta.get("region") or ""), str(meta.get("period") or "")
    ] if b)
    _textbox(s, Inches(0.85), Inches(3.2), Inches(11.6), Inches(1.4),
             [(str(meta.get("name") or "광고 성과 보고서"), 44, True, _INK, PP_ALIGN.LEFT)])
    _textbox(s, Inches(0.9), Inches(4.5), Inches(11.6), Inches(0.6),
             [(sub, 18, True, _MUTED, PP_ALIGN.LEFT)])
    if generated_at:
        _textbox(s, Inches(0.9), Inches(5.2), Inches(6), Inches(0.4),
                 [(generated_at, 12, False, _MUTED, PP_ALIGN.LEFT)])

    # 2) 핵심 지표 (3x2 카드)
    s = _blank(prs)
    _title(s, "핵심 지표")
    cards = [
        ("총 광고비", _won(kpi.get("total_cost", 0)), ""),
        ("노출", _num(kpi.get("total_impressions", 0)), ""),
        ("클릭", _num(kpi.get("total_clicks", 0)), f"CTR {_pct(kpi.get('ctr', 0))} · CPC {_won(kpi.get('cpc', 0))}"),
        ("문의", _num(kpi.get("total_inquiries", 0)), f"1건당 {_won(kpi.get('cpa', 0))}" if kpi.get("total_inquiries") else ""),
        ("단골", _num(kpi.get("total_regulars", 0)), f"1명당 {_won(kpi.get('cpr', 0))}" if kpi.get("total_regulars") else ""),
        ("쿠폰", _num(kpi.get("total_coupons", 0)), f"1건당 {_won(kpi.get('cp_coupon', 0))}" if kpi.get("total_coupons") else ""),
    ]
    cw, ch = Inches(3.85), Inches(2.0)
    gx, gy = Inches(0.6), Inches(1.6)
    for i, (label, value, subv) in enumerate(cards):
        r, c = divmod(i, 3)
        x = Inches(0.6 + c * 4.05)
        y = Inches(1.6 + r * 2.25)
        _card(s, x, y, cw, ch, fill=_WHITE, top_accent=True)
        lines = [(label, 14, True, _MUTED, PP_ALIGN.LEFT), (value, 32, True, _INK, PP_ALIGN.LEFT)]
        if subv:
            lines.append((subv, 12, False, _MUTED, PP_ALIGN.LEFT))
        _textbox(s, x + Inches(0.25), y + Inches(0.25), cw - Inches(0.5), ch - Inches(0.4), lines)

    # 3) 퍼널
    if funnel_stages:
        s = _blank(prs)
        _title(s, "광고 퍼널 — 노출에서 행동까지")
        y = Inches(1.55)
        for st in funnel_stages:
            _card(s, Inches(0.6), y, Inches(12.1), Inches(0.82), fill=_BG)
            _textbox(s, Inches(0.85), y + Inches(0.13), Inches(5), Inches(0.6),
                     [(str(st.get("label", "")), 18, True, _INK, PP_ALIGN.LEFT)])
            meta_bits = []
            if st.get("rate_label") is not None:
                meta_bits.append(f"{st.get('rate_label','')} {_pct(st.get('rate', 0))}")
            if st.get("cost_per", 0):
                meta_bits.append(f"{st.get('cost_label','')} {_won(st.get('cost_per', 0))}")
            _textbox(s, Inches(5.0), y + Inches(0.18), Inches(4.5), Inches(0.5),
                     [(" · ".join(meta_bits), 13, False, _MUTED, PP_ALIGN.LEFT)])
            _textbox(s, Inches(9.5), y + Inches(0.08), Inches(3.0), Inches(0.6),
                     [(_num(st.get("count", 0)), 24, True, _INK, PP_ALIGN.RIGHT)])
            y += Inches(0.95)

    # 4) 한눈에 진단
    conclusion = str(insights.get("conclusion") or "")
    good = str(insights.get("good") or "")
    blocked = str(insights.get("blocked") or "")
    if conclusion or good or blocked:
        s = _blank(prs)
        _title(s, "한눈에 진단")
        if conclusion:
            _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2),
                     [(conclusion, 20, True, _INK, PP_ALIGN.LEFT)])
        gy = Inches(3.0)
        if good:
            _card(s, Inches(0.6), gy, Inches(5.95), Inches(2.4), fill=RGBColor(0xEE, 0xF7, 0xF1), line=RGBColor(0xBF, 0xE3, 0xCC))
            _textbox(s, Inches(0.85), gy + Inches(0.2), Inches(5.5), Inches(2.0),
                     [("잘 된 것", 18, True, RGBColor(0x2E, 0x7D, 0x52), PP_ALIGN.LEFT), (good, 15, False, _INK, PP_ALIGN.LEFT)])
        if blocked:
            _card(s, Inches(6.75), gy, Inches(5.95), Inches(2.4), fill=RGBColor(0xFB, 0xED, 0xE8), line=RGBColor(0xF2, 0xC9, 0xB8))
            _textbox(s, Inches(7.0), gy + Inches(0.2), Inches(5.5), Inches(2.0),
                     [("막힌 것", 18, True, RGBColor(0xB0, 0x40, 0x2A), PP_ALIGN.LEFT), (blocked, 15, False, _INK, PP_ALIGN.LEFT)])

    # 5) 세그먼트 심화
    if segment_rows:
        s = _blank(prs)
        _title(s, "세그먼트 심화 — 어디에 집중할까")
        rows = [[str(r.get("label", "")), _won(r.get("cost", 0)), _won(r.get("cpa", 0)), str(r.get("verdict", ""))]
                for r in segment_rows]
        _table(s, Inches(0.6), Inches(1.6), Inches(12.1),
               ["세그먼트", "비용", "행동당 비용", "판정"], rows,
               col_widths=[Inches(5.1), Inches(2.5), Inches(2.5), Inches(2.0)])

    # 6) 액션
    next_actions = [a for a in (insights.get("next_actions") or []) if str(a).strip()]
    experiments = insights.get("experiments") or []
    if next_actions or experiments:
        s = _blank(prs)
        _title(s, "다음 기간, 이렇게 하겠습니다")
        y = Inches(1.5)
        for i, act in enumerate(next_actions, 1):
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(0.4), Inches(0.4))
            badge.fill.solid(); badge.fill.fore_color.rgb = _ORANGE
            badge.line.fill.background(); badge.shadow.inherit = False
            btf = badge.text_frame; bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run(); br.text = str(i); _set_font(br, 13, True, _WHITE)
            _textbox(s, Inches(1.2), y, Inches(11.4), Inches(0.5),
                     [(str(act), 16, False, _INK, PP_ALIGN.LEFT)])
            y += Inches(0.62)
        if experiments:
            rows = [[str(e.get("priority", "-")), str(e.get("change", "-")),
                     str(e.get("success_criteria", "-")), str(e.get("schedule", "-"))] for e in experiments]
            _table(s, Inches(0.6), y + Inches(0.15), Inches(12.1),
                   ["순위", "변경 내용", "성공 기준", "일정"], rows,
                   col_widths=[Inches(1.1), Inches(5.5), Inches(3.5), Inches(2.0)])

    # 7) 판단 기준
    j = insights.get("judgment") or {}
    if isinstance(j, dict) and (j.get("expand") or j.get("review") or j.get("stop")):
        s = _blank(prs)
        _title(s, "판단 기준")
        specs = [("확대", j.get("expand"), RGBColor(0xEE, 0xF7, 0xF1), RGBColor(0x2E, 0x7D, 0x52)),
                 ("검토", j.get("review"), RGBColor(0xFF, 0xF7, 0xEC), RGBColor(0xB2, 0x6A, 0x00)),
                 ("중단", j.get("stop"), RGBColor(0xFB, 0xED, 0xE8), RGBColor(0xB0, 0x40, 0x2A))]
        x = Inches(0.6)
        for label, text, fill, head in specs:
            if not text:
                continue
            _card(s, x, Inches(1.7), Inches(3.95), Inches(3.4), fill=fill)
            _textbox(s, x + Inches(0.25), Inches(1.95), Inches(3.5), Inches(3.0),
                     [(label, 20, True, head, PP_ALIGN.LEFT), (str(text), 14, False, _INK, PP_ALIGN.LEFT)])
            x += Inches(4.1)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
